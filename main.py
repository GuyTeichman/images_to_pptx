import pptx
from pathlib import Path


def dir_to_ppt(pth, img_suffix):
    txt_width = 19.37
    txt_height = 2
    pic_width = 22.34
    pic_height = 16.69
    prs = pptx.Presentation()
    template = prs.slide_layouts[6]
    print(f'Running {pth.name}...')
    for file in pth.iterdir():
        if file.suffix == img_suffix:
            slide = prs.slides.add_slide(template)
            text = slide.shapes.add_textbox(0, prs.slide_height - txt_height * 360000, pptx.util.Cm(txt_width),
                                            pptx.util.Cm(txt_height))
            text.text = file.name[:-len(img_suffix)]
            pic = slide.shapes.add_picture(str(file), 0, 0, pptx.util.Cm(pic_width), pptx.util.Cm(pic_height))
            pic.crop_bottom, pic.crop_top, pic.crop_left, pic.crop_right = 0, 0, 0, 0

    print(f'Saving {pth.name}...')
    prs.save(pth.joinpath(pth.name + '.pptx'))
    print('Done')


def run(pth='input', image_format='.tif'):
    """
    Create a powerpoint presentation from all of the images in the directory 'pth'. \
    Saves the presentation in the pth, with the name '{directory_names}.pptx'.

    :type pth: str or pathlib.Path
    :param pth: The path for the folder that contains the images. Does not look into subfolders.
    :type image_format: str (default '.tif')
    :param image_format: The expected format of the images ('.jpg','.tif', etc).
    """
    assert isinstance(pth, (str, Path)), "'pth' must be either a string or a pathlib.Path object!"
    if pth == 'input':
        dir_to_ppt(Path(input('Insert the path of the folder containing the images')), img_suffix=image_format)
    else:
        dir_to_ppt(pth=Path(pth), img_suffix=image_format)


def run_subfolders(pth='input', image_format='.tif'):
    """
    Runs 'run()' on every single subfolder within 'pth'.

    :type pth: str or pathlib.Path
    :param pth: The path for the folder that subfolders, that contain images. Only looks one level in.
    :type image_format: str (default '.tif')
    :param image_format: The expected format of the images ('.jpg','.tif', etc).
    """
    assert isinstance(pth, (str, Path)), "'pth' must be either a string or a pathlib.Path object!"
    if pth == 'input':
        pth = Path(input('Insert the path of the folder containing the images'))
    else:
        pth = Path(pth)
    for item in pth.iterdir():
        if item.is_dir():
            run(item, image_format=image_format)


if __name__ == '__main__':
    one_or_multiple = input("Do you want to run a single directory ('single'), or all subfolders ('all')?\n")
    one_or_multiple = one_or_multiple.lower()
    assert one_or_multiple in ['single', 'all'], "Must choose either 'single' or 'all'!"
    pth = input('Insert your full path: \n')
    format = input("What image format do you use? (for example: '.tif', '.jpg', '.png')\n")
    assert format in ['.jpg', '.png', '.tif', '.bmp', '.gif'], f"Invalid image format '{format}'!"
    if one_or_multiple == 'single':
        run(pth=pth, image_format=format)
    else:
        run_subfolders(pth=pth, image_format=format)
